from pptx import Presentation

PATH = r"G:\workburddy\workspace\MyAgent-main\ppt_competition\MyAgent_Competition.pptx"

repl = {
    "[Your Name]": "bry9107795553",
    "[email protected]": "118060862@qq.com",
    "[your-org]": "bry9107795553",
}

p = Presentation(PATH)
replaced = 0
for slide in p.slides:
    for shape in slide.shapes:
        frames = []
        if shape.has_text_frame:
            frames.append(shape.text_frame)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    frames.append(cell.text_frame)
        for tf in frames:
            for para in tf.paragraphs:
                for run in para.runs:
                    for k, v in repl.items():
                        if k in run.text:
                            run.text = run.text.replace(k, v)
                            replaced += 1

p.save(PATH)

# verification: scan for any leftover placeholders
leftover = []
for slide in p.slides:
    for shape in slide.shapes:
        frames = []
        if shape.has_text_frame:
            frames.append(shape.text_frame)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    frames.append(cell.text_frame)
        for tf in frames:
            for para in tf.paragraphs:
                for run in para.runs:
                    for k in repl:
                        if k in run.text:
                            leftover.append(run.text)

print("replaced runs:", replaced)
print("leftover placeholders:", leftover if leftover else "none")
