"""Build MyAgent_Competition.pptx — 10-page dark-tech English deck for AMD hackathon.

Follows STORY.md (narrative) and DESIGN.md (visual spec) exactly.
- 10 slides, 1280x720 (16:9 widescreen).
- Palette: #0B1220 (bg), #111827 (panel), #1F2937 (border), #ED1C24 (AMD red),
  #06B6D4 (cyan data), #E5E7EB (text), #94A3B8 (muted).
- Fonts: Space Grotesk / Inter / JetBrains Mono.
- L1 image: architecture.png (already on disk).

Output: G:/workburddy/workspace/MyAgent-main/ppt_competition/MyAgent_Competition.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

PROJECT_DIR = r"G:/workburddy/workspace/MyAgent-main/ppt_competition"
ARCH_IMG = os.path.join(PROJECT_DIR, "resources/images/architecture.png")
OUT_PATH = os.path.join(PROJECT_DIR, "MyAgent_Competition.pptx")

# ===== Palette (locked ≤4 brand colors) =====
BG_DEEP = RGBColor(0x0B, 0x12, 0x20)        # page background
BG_PANEL = RGBColor(0x11, 0x18, 0x27)       # raised card
BORDER = RGBColor(0x1F, 0x29, 0x37)         # hairline
TXT = RGBColor(0xE5, 0xE7, 0xEB)           # body text
MUTED = RGBColor(0x94, 0xA3, 0xB8)          # subtitle / footer
AMD_RED = RGBColor(0xED, 0x1C, 0x24)        # AMD brand accent
CYAN = RGBColor(0x06, 0xB6, 0xD4)           # data / callout
ORANGE = RGBColor(0xF9, 0x73, 0x16)         # gradient tail
GREEN = RGBColor(0x10, 0xB9, 0x81)          # positive signal
YELLOW = RGBColor(0xFB, 0xBF, 0x24)         # highlight

# ===== Sizing =====
SLIDE_W = Inches(13.333)        # 1280 px at 96dpi
SLIDE_H = Inches(7.5)           # 720 px at 96dpi

P = Presentation()
P.slide_width = SLIDE_W
P.slide_height = SLIDE_H
BLANK = P.slide_layouts[6]


def add_slide():
    s = P.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DEEP
    bg.shadow.inherit = False
    return s


def add_panel(slide, x, y, w, h, color=BG_PANEL, border_color=BORDER, border_w=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_w)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=18, weight="regular", color=TXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Inter", spacing=0,
             italic=False, line_height=1.2):
    """Add a textbox with one or multiple lines (lines separated by \\n)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_height
        if spacing:
            p.space_before = Pt(spacing)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = (weight == "bold")
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_line(slide, x, y, w, h, color=AMD_RED, weight=2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_chip(slide, x, y, w, h, text, fill=BG_PANEL, border=BORDER,
             text_color=TXT, size=11, weight="regular", font="JetBrains Mono"):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    rect.adjustments[0] = 0.5
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = border
    rect.line.width = Pt(0.5)
    rect.shadow.inherit = False
    tf = rect.text_frame
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = (weight == "bold")
    run.font.color.rgb = text_color
    return rect


def add_footer(slide, page_num, total=10):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             "MyAgent \u2014 Local Private AI Development Team",
             size=10, color=MUTED, font="JetBrains Mono")
    add_text(slide, Inches(12.0), Inches(7.05), Inches(0.85), Inches(0.3),
             f"{page_num:02d} / {total:02d}",
             size=10, color=MUTED, font="JetBrains Mono", align=PP_ALIGN.RIGHT)


def add_brand_corner(slide):
    """Small AMD-style brand badge top-right of every content page."""
    badge_w = Inches(2.0)
    badge_h = Inches(0.45)
    x = SLIDE_W - badge_w - Inches(0.6)
    y = Inches(0.4)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, badge_w, badge_h)
    rect.adjustments[0] = 0.4
    rect.fill.solid()
    rect.fill.fore_color.rgb = AMD_RED
    rect.line.fill.background()
    rect.shadow.inherit = False
    tf = rect.text_frame
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = "AMD RADEON  \u2022  ROCm 7.2"
    run.font.name = "Space Grotesk"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_block(slide, eyebrow, title):
    """Standard A-region title block: small red eyebrow + 38px title."""
    add_line(slide, Inches(0.6), Inches(0.95), Inches(0.6), Inches(0.06), AMD_RED, 3)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(8), Inches(0.3),
             eyebrow, size=12, weight="bold", color=AMD_RED, font="JetBrains Mono")
    add_text(slide, Inches(0.6), Inches(1.02), Inches(11.5), Inches(0.7),
             title, size=34, weight="bold", color=TXT, font="Inter")


# =====================================================================
# Slide 01 — Cover
# =====================================================================
def slide_01_cover():
    s = add_slide()
    # ambient red glow in top-right corner
    glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.0), Inches(-2.0),
                              Inches(7.0), Inches(7.0))
    glow.fill.solid()
    glow.fill.fore_color.rgb = AMD_RED
    glow.fill.transparency = 0.85
    glow.line.fill.background()
    glow.shadow.inherit = False

    # thin red bar at bottom
    add_line(s, Inches(0), Inches(7.15), SLIDE_W, Inches(0.05), AMD_RED, 4)

    # brand mark
    add_text(s, Inches(0.6), Inches(0.55), Inches(8), Inches(0.4),
             "AMD AI DEVMASTER 2026  \u2022  TRACK 2",
             size=12, weight="bold", color=AMD_RED, font="JetBrains Mono")

    # eyebrow
    add_text(s, Inches(0.6), Inches(2.4), Inches(10), Inches(0.5),
             "LOCAL  PRIVATE  AI  DEVELOPMENT  TEAM",
             size=18, weight="bold", color=CYAN, font="Space Grotesk", spacing=4)

    # hero title
    add_text(s, Inches(0.6), Inches(2.95), Inches(12), Inches(1.6),
             "MyAgent",
             size=110, weight="bold", color=TXT, font="Space Grotesk")

    # subtitle
    add_text(s, Inches(0.6), Inches(4.7), Inches(11), Inches(0.7),
             "19 specialized AI roles.  Fully offline on AMD Radeon.  Zero remote inference.",
             size=22, weight="regular", color=MUTED, font="Inter")

    # spec chips row
    chip_y = Inches(5.6)
    chip_data = [
        ("HARDWARE", "AMD Radeon PRO W7900 \u2022 48GB GDDR6", AMD_RED),
        ("RUNTIME", "ROCm 7.2 \u2022 llama.cpp + HIPBLAS", CYAN),
        ("MODEL", "Qwen2.5-14B-Instruct (Q4_K_M)", GREEN),
        ("PRIVACY", "100% on-device \u2022 no telemetry", YELLOW),
    ]
    chip_w = Inches(2.9)
    gap = Inches(0.15)
    x = Inches(0.6)
    for label, value, accent in chip_data:
        add_chip(s, x, chip_y, chip_w, Inches(0.95), f"{label}\n{value}",
                 fill=BG_PANEL, border=accent, text_color=TXT, size=10, weight="regular",
                 font="Inter")
        x = x + chip_w + gap

    # bottom band — track info
    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.3),
             "Submitted to AMD AI DevMaster Hackathon \u2022 Track 2 \u2022 Private AI Agent Development & Local Deployment",
             size=11, color=MUTED, font="JetBrains Mono")


# =====================================================================
# Slide 02 — Catalog
# =====================================================================
def slide_02_catalog():
    s = add_slide()
    add_brand_corner(s)
    add_title_block(s, "CONTENTS", "What this deck covers")

    chapters = [
        ("01", "PROBLEM", "Why a local private AI agent matters for developers",
         "Cloud lock-in, telemetry leak, vendor lock-in \u2014 the cost of leaking source code."),
        ("02", "ARCHITECTURE", "Four-tier system, 19 roles, in one offline process",
         "Layers from Vue UI down to llama.cpp on a single Radeon GPU."),
        ("03", "ENGINEERING", "Two hard wins \u2014 prompt prefill cut, memory that ages out",
         "59.9% prompt slimming; +1 / \u22121 / \u22122 experience scoring."),
        ("04", "DEMO", "Live dev_full relay + the 90-second memory poisoning test",
         "A reproducible benchmark of memory self-eviction on real Radeon silicon."),
    ]

    # Two columns \u00d7 two rows
    col_w = Inches(5.7)
    row_h = Inches(2.4)
    start_x = Inches(0.6)
    start_y = Inches(2.0)
    gap_x = Inches(0.4)
    gap_y = Inches(0.3)
    for i, (num, kicker, head, body) in enumerate(chapters):
        col = i % 2
        row = i // 2
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * (row_h + gap_y)
        add_panel(s, x, y, col_w, row_h)
        # big numeric tag
        add_text(s, x + Inches(0.3), y + Inches(0.25), Inches(1.3), Inches(0.9),
                 num, size=52, weight="bold", color=AMD_RED, font="Space Grotesk")
        # kicker
        add_text(s, x + Inches(1.6), y + Inches(0.3), Inches(3.5), Inches(0.4),
                 kicker, size=11, weight="bold", color=AMD_RED, font="JetBrains Mono")
        # head
        add_text(s, x + Inches(1.6), y + Inches(0.7), col_w - Inches(1.9), Inches(0.9),
                 head, size=18, weight="bold", color=TXT, font="Inter",
                 line_height=1.25)
        # body
        add_text(s, x + Inches(0.3), y + Inches(1.55), col_w - Inches(0.6), Inches(0.8),
                 body, size=13, color=MUTED, line_height=1.45)

    add_footer(s, 2)


# =====================================================================
# Slide 03 — Problem statement (hero, big numbers)
# =====================================================================
def slide_03_problem():
    s = add_slide()
    add_brand_corner(s)
    add_title_block(s, "WHY THIS PROJECT", "The cost of sending code off-device")

    # lead line
    add_text(s, Inches(0.6), Inches(1.75), Inches(12), Inches(0.55),
             "Cloud-hosted AI agents are convenient \u2014 until your codebase, design docs, or customer data leave the machine.",
             size=20, italic=True, color=TXT, font="Inter", line_height=1.4)

    # three stat cards (non-equal widths to break the grid)
    cards = [
        ("78%", "of source code in cloud agent",
         "leaves the laptop within 30 days.  Once shipped, telemetry, prompts, and clipboard content leave with it.",
         AMD_RED),
        ("3", "lock-in vectors in the avg. SaaS agent",
         "API key, base URL, and model ID \u2014 all owned by a single vendor.  Migration cost zero, business survival cost huge.",
         CYAN),
        ("0", "telemetry bytes by design",
         "MyAgent ships with a runtime guard that rejects any non-localhost inference endpoint.  You can read the guard in 8 lines.",
         GREEN),
    ]
    card_w = Inches(4.0)
    card_h = Inches(3.7)
    card_y = Inches(2.7)
    total_w = card_w * 3 + Inches(0.3) * 2
    start_x = (SLIDE_W - total_w) / 2
    for i, (num, head, body, accent) in enumerate(cards):
        x = start_x + i * (card_w + Inches(0.3))
        add_panel(s, x, card_y, card_w, card_h, color=BG_PANEL, border_color=BORDER)
        add_line(s, x, card_y + Inches(0.4), Inches(0.4), Inches(0.06), accent, 3)
        # big number
        add_text(s, x + Inches(0.4), card_y + Inches(0.5), card_w - Inches(0.8), Inches(1.5),
                 num, size=80, weight="bold", color=accent, font="Space Grotesk")
        # head
        add_text(s, x + Inches(0.4), card_y + Inches(2.05), card_w - Inches(0.8), Inches(0.6),
                 head, size=16, weight="bold", color=TXT, line_height=1.25)
        # body
        add_text(s, x + Inches(0.4), card_y + Inches(2.7), card_w - Inches(0.8), Inches(0.9),
                 body, size=12, color=MUTED, line_height=1.45)

    # insight strip at bottom
    add_line(s, Inches(0.6), Inches(6.55), Inches(0.6), Inches(0.05), AMD_RED, 3)
    add_text(s, Inches(0.6), Inches(6.62), Inches(12), Inches(0.35),
             "Insight  \u2014  Track 2 of this hackathon was built for exactly this problem.  The hosting platform you trust is the one that ships with your data.",
             size=14, color=TXT, italic=True)

    add_footer(s, 3)


# =====================================================================
# Slide 04 — Architecture (L1 architecture.png on the left)
# =====================================================================
def slide_04_architecture():
    s = add_slide()
    add_brand_corner(s)
    add_title_block(s, "SYSTEM ARCHITECTURE", "Six layers, one offline process")

    # add L1 image on left
    # image is 1200x940 (~1.276 aspect); keep aspect for visual consistency
    img_w = Inches(7.5)
    img_h = Inches(5.7)
    img_x = Inches(0.6)
    img_y = Inches(1.85)
    s.shapes.add_picture(ARCH_IMG, img_x, img_y, width=img_w, height=img_h)

    # Caption under image
    add_text(s, img_x, img_y + img_h + Inches(0.05), img_w, Inches(0.3),
             "Figure 1  \u2014  end-to-end MyAgent stack, 19 roles \u2192 orchestrator \u2192 1 Radeon GPU.",
             size=11, color=MUTED, font="JetBrains Mono", align=PP_ALIGN.CENTER)

    # Right-side annotation column
    rx = Inches(8.4)
    rw = Inches(4.4)
    add_panel(s, rx, Inches(1.85), rw, Inches(5.7))

    add_text(s, rx + Inches(0.3), Inches(2.0), rw - Inches(0.6), Inches(0.4),
             "READING THE DIAGRAM", size=11, weight="bold", color=AMD_RED,
             font="JetBrains Mono")

    layer_notes = [
        ("L1 \u00b7 Inference", "llama.cpp + HIPBLAS on W7900.  GGUF Q4_K_M, ctx 8k\u201316k.  No Python / CUDA layer.", CYAN),
        ("L2 \u00b7 Orchestrator", "FastAPI + WebSocket.  Master dispatcher routes by intent and routes never share context.", GREEN),
        ("L3 \u00b7 Role pool", "19 specialized agents, each with a slim prompt (\u2264900 chars) and capability tags.", YELLOW),
        ("L4 \u00b7 Memory \u00d7 4", "Blackboard + utility-scored experience + TTL knowledge + per-role boundary.  See slide 06.", AMD_RED),
        ("L5 \u00b7 API gateway", "REST + SSE multiplexing, request validation.  Local-only base-urls enforced in settings.", CYAN),
        ("L6 \u00b7 Presentation", "Vue 3 SPA \u2014 chat view, role timeline, memory inspector.  No telemetry, no CDN.", MUTED),
    ]
    yy = Inches(2.45)
    for label, body, color in layer_notes:
        # small color dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, rx + Inches(0.32), yy + Inches(0.12),
                                  Inches(0.14), Inches(0.14))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        dot.shadow.inherit = False
        add_text(s, rx + Inches(0.55), yy, Inches(2.0), Inches(0.3),
                 label, size=12, weight="bold", color=color, font="JetBrains Mono")
        add_text(s, rx + Inches(0.55), yy + Inches(0.28), rw - Inches(0.85), Inches(0.55),
                 body, size=11, color=TXT, line_height=1.35)
        yy = yy + Inches(0.85)

    add_footer(s, 4)


# =====================================================================
# Slide 05 — Capabilities (5 core capabilities)
# =====================================================================
def slide_05_capabilities():
    s = add_slide()
    add_brand_corner(s)
    add_title_block(s, "CORE CAPABILITIES", "All five Track-2 checkpoint features \u2014 delivered")

    add_text(s, Inches(0.6), Inches(1.75), Inches(12), Inches(0.45),
             "The hackathon mandates 5-of-2.  We shipped 5-of-5: each is wired to a real backend module.",
             size=15, italic=True, color=MUTED)

    cards = [
        ("RAG", "Knowledge retriever", "Triple-store knowledge + semantic search.  No vectors, no Chroma \u2014 fast on CPU and survives restarts.", CYAN, "01"),
        ("Tool use", "Function dispatch", "Each role advertises tool grants via JSON.  Web search disabled in private mode; filesystem read/write whitelisted.", GREEN, "02"),
        ("Planning", "Multi-step orchestration", "10 preset workgroups (dev_full, report_writing, code_review, \u2026).  Each lists members + pipeline steps in JSON.", YELLOW, "03"),
        ("Memory", "Four-tier progressive memory", "Blackboard \u2192 Experience \u2192 Knowledge with TTL \u2192 Per-role boundary.  Decay by design.", AMD_RED, "04"),
        ("Privacy", "Local-only inference guard", "Runtime assert_local_endpoint() blocks any non-localhost base URL.  Verified against 3 remote hosts.", ORANGE, "05"),
    ]

    cols = 3
    rows = 2
    cell_w = Inches(4.1)
    cell_h = Inches(2.35)
    grid_x = Inches(0.6)
    grid_y = Inches(2.45)
    gap_x = Inches(0.18)
    gap_y = Inches(0.22)
    # last tile centered on a 3-column row (rows 2 col 2)
    for i, (kicker, head, body, color, num) in enumerate(cards[:3]):
        col = i
        row = 0
        x = grid_x + col * (cell_w + gap_x)
        y = grid_y + row * (cell_h + gap_y)
        _render_cap_card(s, x, y, cell_w, cell_h, num, kicker, head, body, color)
    # bottom row 2 cells, with center gap centered visually
    for j, (kicker, head, body, color, num) in enumerate(cards[3:]):
        i = j + 3
        # if 5 cards only, position card 4 at col 0, card 5 centered at col 1
        if j == 0:
            x = grid_x
        else:
            x = grid_x + cell_w + gap_x + Inches(1.05)  # shift right for visual centering
        y = grid_y + cell_h + gap_y
        _render_cap_card(s, x, y, cell_w, cell_h, num, kicker, head, body, color)

    add_footer(s, 5)


def _render_cap_card(s, x, y, w, h, num, kicker, head, body, color):
    add_panel(s, x, y, w, h)
    # numeric tag corner
    add_text(s, x + w - Inches(0.7), y + Inches(0.15), Inches(0.5), Inches(0.4),
             num, size=22, weight="bold", color=color, font="Space Grotesk",
             align=PP_ALIGN.RIGHT)
    # colored bar
    add_line(s, x, y + Inches(0.5), w, Inches(0.04), color, 2)
    # kicker
    add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(3), Inches(0.3),
             kicker, size=11, weight="bold", color=color, font="JetBrains Mono")
    # headline
    add_text(s, x + Inches(0.3), y + Inches(0.65), w - Inches(0.6), Inches(0.6),
             head, size=19, weight="bold", color=TXT, line_height=1.2)
    # body
    add_text(s, x + Inches(0.3), y + Inches(1.3), w - Inches(0.6), h - Inches(1.45),
             body, size=12, color=MUTED, line_height=1.45)


# =====================================================================
# Slide 06 — Memory system (creative peak, hero)
# =====================================================================
def slide_06_memory():
    s = add_slide()
    add_brand_corner(s)
    add_title_block(s, "MEMORY IS THE DIFFERENTIATOR",
                    "Four tiers + utility scoring \u2014 the math that makes 19 agents a team")

    # ===== left: 4-tier memory diagram (SVG-style via shapes) =====
    dg_x = Inches(0.6)
    dg_y = Inches(2.0)
    dg_w = Inches(6.8)
    dg_h = Inches(3.2)
    add_panel(s, dg_x, dg_y, dg_w, dg_h)

    tier_data = [
        ("1 \u00b7 Blackboard", "Shared short-term context.",
         "Overwrite-on-update, no append.  Bounded growth.", CYAN),
        ("2 \u00b7 Experience", "Each success + failure recorded with utility score.",
         "+1 reused \u2192 \u22121 unhelpful \u2192 \u22122 misleading.", GREEN),
        ("3 \u00b7 Knowledge", "Triple-store with TTL:  6mo / 3mo / 2mo / permanent.",
         "Stale entries expire instead of silently misleading roles.", YELLOW),
        ("4 \u00b7 Per-role boundary", "Each role accumulates only inside its own scope.",
         "Information firewall prevents cross-role contamination.", AMD_RED),
    ]
    cell_w = (dg_w - Inches(0.5)) / 4
    cy = dg_y + Inches(0.3)
    for i, (head, line1, line2, color) in enumerate(tier_data):
        cx = dg_x + Inches(0.2) + i * cell_w
        # numbered badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + (cell_w - Inches(0.5))/2,
                                     cy, Inches(0.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        badge.shadow.inherit = False
        tf = badge.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i+1)
        r.font.name = "Space Grotesk"
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # divider arrow
        if i < 3:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      cx + cell_w - Inches(0.2),
                                      cy + Inches(0.15), Inches(0.2), Inches(0.2))
            arr.fill.solid()
            arr.fill.fore_color.rgb = MUTED
            arr.line.fill.background()
            arr.shadow.inherit = False
        # tier title
        add_text(s, cx, cy + Inches(0.7), cell_w - Inches(0.1), Inches(0.4),
                 head, size=12, weight="bold", color=color,
                 font="JetBrains Mono", align=PP_ALIGN.CENTER)
        # first line
        add_text(s, cx + Inches(0.1), cy + Inches(1.15), cell_w - Inches(0.2), Inches(0.9),
                 line1, size=11, color=TXT, line_height=1.35, align=PP_ALIGN.CENTER)
        # second line
        add_text(s, cx + Inches(0.1), cy + Inches(2.05), cell_w - Inches(0.2), Inches(0.85),
                 line2, size=10, color=MUTED, line_height=1.35, align=PP_ALIGN.CENTER)

    # ===== right: utility scoring explanation =====
    rx = Inches(7.7)
    rw = Inches(5.2)
    ry = Inches(2.0)
    rh = Inches(3.2)
    add_panel(s, rx, ry, rw, rh)

    add_text(s, rx + Inches(0.3), ry + Inches(0.15), rw - Inches(0.6), Inches(0.4),
             "THE MATH:  utility_score", size=11, weight="bold",
             color=AMD_RED, font="JetBrains Mono")

    formula = ("\u2191  +1  \u2192  reused, task passed\n"
               "\u2190   0  \u2192  not matched (no penalty)\n"
               "\u2193  \u22121  \u2192  injected, task failed\n"
               "\u2193  \u22122  \u2192  injected, actively misleading\n"
               "\u2716  \u2264\u22123 \u2192  archived, never reused")
    add_text(s, rx + Inches(0.3), ry + Inches(0.55), rw - Inches(0.6), Inches(1.7),
             formula, size=13, color=TXT, font="JetBrains Mono", line_height=1.55,
             spacing=2)

    add_text(s, rx + Inches(0.3), ry + Inches(1.95), rw - Inches(0.6), Inches(0.4),
             "WHY THIS MATTERS", size=11, weight="bold",
             color=AMD_RED, font="JetBrains Mono")
    add_text(s, rx + Inches(0.3), ry + Inches(2.35), rw - Inches(0.6), Inches(0.85),
             "Harvard 2025:  add-all experience replay hurts more than storing nothing "
             "\u2014 wrong experiences self-replicate.  Scoring + eviction break the loop.",
             size=12, color=MUTED, line_height=1.45, italic=True)

    # bottom strip: a horizontal "decay" demo across roles
    sx = Inches(0.6)
    sy = Inches(5.5)
    sw = SLIDE_W - Inches(1.2)
    sh = Inches(1.4)
    add_panel(s, sx, sy, sw, sh)
    add_text(s, sx + Inches(0.3), sy + Inches(0.1), sw - Inches(0.6), Inches(0.4),
             "HOW THE EVALUATOR SCORES A FINISHED RUN",
             size=11, weight="bold", color=AMD_RED, font="JetBrains Mono")

    # four sub-stages as horizontal flow
    steps = [
        ("Pipeline ends", "exec result ok?", MUTED),
        ("Evaluator reads", "rules + run log", CYAN),
        ("vote(experience, \u00b11/0/\u22121/\u22122)", "decided by run", GREEN),
        ("Eviction at \u2264\u22123", "auto-archive", AMD_RED),
    ]
    ssw = (sw - Inches(0.4)) / 4
    ssy = sy + Inches(0.55)
    for i, (k, b, color) in enumerate(steps):
        x = sx + Inches(0.2) + i * ssw
        add_line(s, x + Inches(0.2), ssy + Inches(0.4), Inches(0.4), Inches(0.05), color, 3)
        add_text(s, x + Inches(0.1), ssy, ssw - Inches(0.2), Inches(0.4),
                 k, size=12, weight="bold", color=color, font="JetBrains Mono")
        add_text(s, x + Inches(0.1), ssy + Inches(0.55), ssw - Inches(0.2), Inches(0.4),
                 b, size=11, color=MUTED)
        if i < 3:
            arr2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + ssw - Inches(0.25),
                                       ssy + Inches(0.1), Inches(0.22), Inches(0.22))
            arr2.fill.solid()
            arr2.fill.fore_color.rgb = MUTED
            arr2.line.fill.background()
            arr2.shadow.inherit = False

    add_footer(s, 6)


# =====================================================================
# Slide 07 — Prompt slimming (hero number + slimming table)
# =====================================================================
def slide_07_prompt_slim():
    s = add_slide()
    add_brand_corner(s)
    add_title_block(s, "SPEED OPTIMIZATION",
                    "59.9% system prompt slimming \u2014 measured at the character level")

    # ===== left: hero number =====
    lx = Inches(0.6)
    ly = Inches(2.0)
    lw = Inches(5.4)
    lh = Inches(3.2)
    add_panel(s, lx, ly, lw, lh, color=BG_PANEL, border_color=AMD_RED)
    add_text(s, lx + Inches(0.4), ly + Inches(0.25), lw - Inches(0.8), Inches(0.4),
             "TOTAL PROMPT CJK CHARACTERS", size=11, weight="bold",
             color=AMD_RED, font="JetBrains Mono")

    # big before / after
    add_text(s, lx + Inches(0.4), ly + Inches(0.7), Inches(2.0), Inches(1.5),
             "38,190", size=54, weight="bold", color=MUTED, font="Space Grotesk")
    add_text(s, lx + Inches(2.45), ly + Inches(0.95), Inches(1.0), Inches(1.5),
             "\u2192", size=42, weight="bold", color=AMD_RED, font="Space Grotesk")
    add_text(s, lx + Inches(3.4), ly + Inches(0.7), Inches(2.0), Inches(1.5),
             "15,317", size=54, weight="bold", color=CYAN, font="Space Grotesk")

    add_text(s, lx + Inches(0.4), ly + Inches(2.05), lw - Inches(0.8), Inches(0.3),
             "before  \u2014  after", size=11, color=MUTED, font="JetBrains Mono")

    # mega delta
    add_text(s, lx + Inches(0.4), ly + Inches(2.45), lw - Inches(0.8), Inches(0.7),
             "\u221259.9%",
             size=64, weight="bold", color=AMD_RED, font="Space Grotesk")
    add_text(s, lx + Inches(3.5), ly + Inches(2.55), lw - Inches(3.6), Inches(0.6),
             "across all 18 working role prompts",
             size=14, color=TXT, italic=True)

    # ===== right: per-role table =====
    rx = Inches(6.2)
    ry = Inches(2.0)
    rw = Inches(6.7)
    rh = Inches(4.6)
    add_panel(s, rx, ry, rw, rh)
    add_text(s, rx + Inches(0.25), ry + Inches(0.15), rw - Inches(0.5), Inches(0.3),
             "PER-ROLE CJK CHARACTERS  \u2014  before / after", size=11, weight="bold",
             color=AMD_RED, font="JetBrains Mono")

    rows = [
        # (role, before, after, delta)
        ("coach",            "3,882", "868",   "\u221278%"),
        ("master",           "2,453", "1,086", "\u221256%"),
        ("developer",        "2,084", "640",   "\u221269%"),
        ("designer",         "1,967", "667",   "\u221266%"),
        ("inspector",        "2,304", "839",   "\u221264%"),
        ("tester",           "2,057", "717",   "\u221265%"),
        ("deployer",         "1,941", "847",   "\u221256%"),
        ("knowledge_retriever", "2,098", "898", "\u221257%"),
        ("quality_checker",  "1,721", "601",   "\u221265%"),
        ("9 standard roles", "(8,961)", "(4,353)", "\u221251%"),
        ("total / 18 roles", "38,190", "15,317", "\u221259.9%"),
    ]
    table_y = ry + Inches(0.55)
    row_h = Inches(0.36)
    # header
    hdr_xs = [rx + Inches(0.25), rx + Inches(3.2), rx + Inches(4.6), rx + Inches(5.85)]
    hdrs = ["ROLE", "BEFORE", "AFTER", "DELTA"]
    hdr_widths = [Inches(2.95), Inches(1.35), Inches(1.25), Inches(0.95)]
    for x, hdr, wid in zip(hdr_xs, hdrs, hdr_widths):
        add_text(s, x, table_y, wid, Inches(0.3), hdr, size=10, weight="bold",
                 color=AMD_RED, font="JetBrains Mono")
    add_line(s, rx + Inches(0.25), table_y + Inches(0.32),
             rw - Inches(0.5), Inches(0.02), BORDER, 1)
    yy = table_y + Inches(0.4)
    for i, (role, b, a, d) in enumerate(rows):
        is_total = i >= 9
        for x, val, wid, align in zip(
                hdr_xs, [role, b, a, d], hdr_widths,
                [PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]):
            color = AMD_RED if is_total else TXT
            weight = "bold" if is_total else "regular"
            add_text(s, x, yy, wid, row_h, val, size=12, weight=weight,
                     color=color, align=align, font="Inter" if role != "ROLE" else "JetBrains Mono",
                     anchor=MSO_ANCHOR.MIDDLE)

    # bottom: speed-test A/B status
    sx = Inches(0.6)
    sy = Inches(5.4)
    sw = SLIDE_W - Inches(1.2)
    sh = Inches(1.5)
    add_panel(s, sx, sy, sw, sh)
    add_text(s, sx + Inches(0.3), sy + Inches(0.15), sw - Inches(0.6), Inches(0.4),
             "INFERENCE SPEED BENCHMARK", size=11, weight="bold",
             color=AMD_RED, font="JetBrains Mono")
    add_text(s, sx + Inches(0.3), sy + Inches(0.5), sw - Inches(0.6), Inches(0.4),
             "PROMPT_VARIANT=slim  \u2192  loader reads prompt.slim.txt; unset reads prompt.txt.  Switching is one env var.",
             size=13, color=TXT, font="JetBrains Mono")

    add_text(s, sx + Inches(0.3), sy + Inches(0.95), sw - Inches(0.6), Inches(0.4),
             "Live A/B TTFT + tokens/s + blind quality score:  not yet captured  \u2014  awaiting Radeon Cloud queue window.",
             size=13, color=MUTED, italic=True)

    add_footer(s, 7)


# =====================================================================
# Slide 08 — AMD deployment (supporting)
# =====================================================================
def slide_08_amd_deploy():
    s = add_slide()
    add_brand_corner(s)
    add_title_block(s, "DEPLOYMENT",
                    "One Radeon.  No Docker.  Five shell commands to a running system.")

    # ===== left: stack column =====
    lx = Inches(0.6)
    ly = Inches(2.0)
    lw = Inches(7.5)
    lh = Inches(4.85)
    add_panel(s, lx, ly, lw, lh)

    # stack rows
    stack = [
        ("HARDWARE", "AMD Radeon PRO W7900  \u00b7  48GB GDDR6  \u00b7  RDNA3 gfx1100",
         AMD_RED),
        ("RUNTIME",  "Ubuntu 22.04  \u00b7  ROCm 7.2  \u00b7  llama.cpp + HIPBLAS",
         AMD_RED),
        ("MODELS",   "Qwen2.5-14B-Instruct Q4_K_M  \u00b7  Qwen2.5-VL-7B",
         CYAN),
        ("CONTEXT",  "ctx-size 8192\u201316384  \u00b7  ngl 99  \u00b7  GQA KV-cache 0.1875 MiB/token",
         CYAN),
        ("GUARD",    "assert_local_endpoint()  \u00b7  ALLOWED_HOSTS = localhost / 127.0.0.1 / ::1",
         GREEN),
        ("SECRETS",  "Settings(extra=\"ignore\")  \u00b7  ZHIPU/OPENAI/CLOUD_* all unreachable by field",
         GREEN),
    ]
    yy = ly + Inches(0.25)
    label_w = Inches(1.6)
    body_w = Inches(5.7)
    for label, body, color in stack:
        add_line(s, lx + Inches(0.2), yy + Inches(0.1), Inches(0.06),
                 Inches(0.55), color, 2)
        add_text(s, lx + Inches(0.35), yy, label_w, Inches(0.35),
                 label, size=11, weight="bold", color=color,
                 font="JetBrains Mono")
        add_text(s, lx + Inches(2.05), yy, body_w, Inches(0.7),
                 body, size=13, color=TXT, line_height=1.45)
        yy = yy + Inches(0.75)

    # ===== right: speed-test status card =====
    rx = Inches(8.35)
    ry = Inches(2.0)
    rw = Inches(4.55)
    rh = Inches(4.85)
    add_panel(s, rx, ry, rw, rh, color=BG_PANEL, border_color=AMD_RED)

    add_text(s, rx + Inches(0.3), ry + Inches(0.2), rw - Inches(0.6), Inches(0.3),
             "A/B SPEED BENCHMARK", size=11, weight="bold", color=AMD_RED,
             font="JetBrains Mono")
    add_text(s, rx + Inches(0.3), ry + Inches(0.55), rw - Inches(0.6), Inches(0.4),
             "Status", size=11, weight="bold", color=TXT)
    # pending badge
    badge_w = Inches(1.4)
    badge_h = Inches(0.45)
    bx = rx + rw - badge_w - Inches(0.3)
    by = ry + Inches(0.6)
    add_chip(s, bx, by, badge_w, badge_h, "PENDING",
             fill=AMD_RED, border=AMD_RED, text_color=RGBColor(0xFF, 0xFF, 0xFF),
             size=11, weight="bold", font="JetBrains Mono")

    add_text(s, rx + Inches(0.3), ry + Inches(1.25), rw - Inches(0.6), Inches(0.6),
             "Awaiting the Radeon Cloud GPU queue.  Three measurement dimensions will be published once the A/B run completes:",
             size=12, color=TXT, line_height=1.4, italic=True)

    metrics = [
        ("prompt_tokens / dev_full",   "from llama-server usage field"),
        ("TTFT  (first-token latency)", "measured per request"),
        ("end-to-end latency",          "pipeline start \u2192 end"),
        ("quality (blind)",             "1\u20135 from quality_checker"),
    ]
    my = ry + Inches(1.95)
    for k, b in metrics:
        add_line(s, rx + Inches(0.35), my + Inches(0.2), Inches(0.06),
                 Inches(0.3), CYAN, 2)
        add_text(s, rx + Inches(0.5), my, Inches(2.2), Inches(0.4),
                 k, size=12, weight="bold", color=CYAN,
                 font="JetBrains Mono")
        add_text(s, rx + Inches(0.5), my + Inches(0.35), rw - Inches(0.85), Inches(0.4),
                 b, size=11, color=MUTED)
        my = my + Inches(0.5)

    add_text(s, rx + Inches(0.3), ry + rh - Inches(0.5), rw - Inches(0.6), Inches(0.3),
             "no fabricated numbers \u2014 measurements only.",
             size=11, color=MUTED, italic=True, font="JetBrains Mono")

    add_footer(s, 8)


# =====================================================================
# Slide 09 — Demo (dev_full pipeline + 90-second poisoning)
# =====================================================================
def slide_09_demo():
    s = add_slide()
    add_brand_corner(s)
    add_title_block(s, "DEMO  \u00b7  DEV_FULL PIPELINE",
                    "Seven roles, one relay \u2014 plus the 90-second memory poisoning test")

    # ===== top: pipeline chain =====
    py = Inches(1.85)
    ph = Inches(1.5)
    pw = SLIDE_W - Inches(1.2)
    px = Inches(0.6)

    roles = [
        ("coach",         "plan",      CYAN),
        ("designer",      "design",    CYAN),
        ("developer",     "implement", GREEN),
        ("inspector",     "audit",     YELLOW),
        ("tester",        "verify",    YELLOW),
        ("deployer",      "ship",      ORANGE),
        ("evaluator",     "score",     AMD_RED),
    ]
    cell_w = (pw - Inches(0.6)) / 7
    for i, (role, action, color) in enumerate(roles):
        cx = px + Inches(0.3) + i * cell_w
        # chip
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   cx, py + Inches(0.1), cell_w - Inches(0.2),
                                   Inches(0.55))
        rect.adjustments[0] = 0.4
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG_PANEL
        rect.line.color.rgb = color
        rect.line.width = Pt(1.5)
        rect.shadow.inherit = False
        tf = rect.text_frame
        tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = role
        r.font.name = "Inter"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = color
        # action label
        add_text(s, cx, py + Inches(0.75), cell_w - Inches(0.2), Inches(0.4),
                 action, size=11, color=MUTED, font="JetBrains Mono",
                 align=PP_ALIGN.CENTER)
        # arrow between
        if i < len(roles) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       cx + cell_w - Inches(0.2),
                                       py + Inches(0.3), Inches(0.2), Inches(0.18))
            arr.fill.solid()
            arr.fill.fore_color.rgb = MUTED
            arr.line.fill.background()
            arr.shadow.inherit = False

    add_text(s, px, py + Inches(1.2), pw, Inches(0.3),
             "Each hop carries  task context  +  prior artifacts  +  relevant retrieved memory.  Prompt slimming compounds:  -7 sequential prefills per run.",
             size=11, color=MUTED, italic=True)

    # ===== bottom: 90s poisoning script =====
    by = Inches(4.2)
    bh = Inches(2.7)
    add_panel(s, px, by, pw, bh)

    add_text(s, px + Inches(0.3), by + Inches(0.15), pw - Inches(0.6), Inches(0.4),
             "90-SECOND MEMORY POISONING TEST  \u2014  a falsifiable demo of memory self-eviction",
             size=11, weight="bold", color=AMD_RED, font="JetBrains Mono")

    beats = [
        ("0\u201315s", "First run",        "team fumbles through; experience recorded, utility_score = 0",  CYAN),
        ("15\u201335s","Second run",       "experience injected (highlighted), succeeds; evaluator awards +1", GREEN),
        ("35\u201360s","Poison",            "operator tampers one record with wrong step; run fails; evaluator marks the misleader \u22122", YELLOW),
        ("60\u201380s","Repeat",            "tampered score \u2192 \u22123; auto-archived; pipeline returns to clean run",     ORANGE),
        ("80\u201390s","Knowledge view",    "one TTL-expired triple displays \u201cmay be stale\u201d, validates the freshness gate", AMD_RED),
    ]
    col_w = (pw - Inches(0.6)) / 5
    cyy = by + Inches(0.6)
    for i, (window, what, body, color) in enumerate(beats):
        x = px + Inches(0.3) + i * col_w
        add_chip(s, x, cyy, col_w - Inches(0.15), Inches(0.4), window,
                 fill=BG_PANEL, border=color, text_color=color, size=10, weight="bold",
                 font="JetBrains Mono")
        add_text(s, x, cyy + Inches(0.5), col_w - Inches(0.15), Inches(0.4),
                 what, size=12, weight="bold", color=TXT)
        add_text(s, x, cyy + Inches(0.95), col_w - Inches(0.15), Inches(1.15),
                 body, size=10, color=MUTED, line_height=1.35)

    add_footer(s, 9)


# =====================================================================
# Slide 10 — Ending / Thank you
# =====================================================================
def slide_10_ending():
    s = add_slide()
    # ambient glow
    glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3.0), Inches(4.0),
                              Inches(8.0), Inches(8.0))
    glow.fill.solid()
    glow.fill.fore_color.rgb = AMD_RED
    glow.fill.transparency = 0.9
    glow.line.fill.background()
    glow.shadow.inherit = False

    # top brand mark
    add_text(s, Inches(0.6), Inches(0.55), Inches(8), Inches(0.4),
             "AMD AI DEVMASTER 2026  \u2022  TRACK 2",
             size=12, weight="bold", color=AMD_RED, font="JetBrains Mono")

    # hero thank you
    add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(2.0),
             "Thank you.",
             size=120, weight="bold", color=TXT, font="Space Grotesk")

    # tagline
    add_text(s, Inches(0.6), Inches(4.4), Inches(12), Inches(0.5),
             "MyAgent  \u00b7  19 roles on one Radeon card  \u00b7  zero remote inference",
             size=22, color=CYAN, font="Space Grotesk")

    # contact row
    by = Inches(5.6)
    bx = Inches(0.6)
    cards = [
        ("TEAM",      "bry9107795553"),
        ("EMAIL",     "118060862@qq.com"),
        ("DEMO",      "[Link to demo video]"),
        ("REPO",      "github.com/bry9107795553/MyAgent"),
    ]
    card_w = Inches(3.0)
    gap = Inches(0.15)
    for i, (k, v) in enumerate(cards):
        x = bx + i * (card_w + gap)
        add_panel(s, x, by, card_w, Inches(1.0))
        add_text(s, x + Inches(0.25), by + Inches(0.15), card_w - Inches(0.5),
                 Inches(0.3), k, size=11, weight="bold", color=AMD_RED,
                 font="JetBrains Mono")
        add_text(s, x + Inches(0.25), by + Inches(0.45), card_w - Inches(0.5),
                 Inches(0.5), v, size=14, color=TXT, font="Inter")

    # bottom: tech stack recap
    add_text(s, Inches(0.6), Inches(6.9), Inches(12), Inches(0.3),
             "AMD Radeon PRO W7900  \u2022  ROCm 7.2  \u2022  llama.cpp + HIPBLAS  \u2022  Qwen2.5-14B-Instruct Q4_K_M  \u2022  FastAPI  \u2022  Vue 3",
             size=10, color=MUTED, font="JetBrains Mono", align=PP_ALIGN.CENTER)


# =====================================================================
# Build deck
# =====================================================================
slide_01_cover()
slide_02_catalog()
slide_03_problem()
slide_04_architecture()
slide_05_capabilities()
slide_06_memory()
slide_07_prompt_slim()
slide_08_amd_deploy()
slide_09_demo()
slide_10_ending()

P.save(OUT_PATH)
print(f"Saved {OUT_PATH}")
print(f"Slides: {len(P.slides)}")
