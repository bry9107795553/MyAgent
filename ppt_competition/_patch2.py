from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PATH = r"G:\workburddy\workspace\MyAgent-main\ppt_competition\MyAgent_Competition.pptx"
repl = {
    "[Your Name]": "bry9107795553",
    "[email protected]": "118060862@qq.com",
    "[your-org]": "bry9107795553",
}

def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from iter_shapes(sh.shapes)
        except Exception:
            pass

p = Presentation(PATH)
hits = []
for si, slide in enumerate(p.slides):
    for sh in iter_shapes(slide.shapes):
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if any(k in run.text for k in repl):
                    hits.append((si, repr(run.text)))
print("found placeholders:", hits)

replaced = 0
for slide in p.slides:
    for sh in iter_shapes(slide.shapes):
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                for k, v in repl.items():
                    if k in run.text:
                        run.text = run.text.replace(k, v)
                        replaced += 1
p.save(PATH)

leftover = []
for slide in p.slides:
    for sh in iter_shapes(slide.shapes):
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if any(k in run.text for k in repl):
                    leftover.append(run.text)
print("replaced:", replaced, "leftover:", leftover if leftover else "none")
