from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PATH = r"G:\workburddy\workspace\MyAgent-main\ppt_competition\MyAgent_Competition.pptx"

def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from iter_shapes(sh.shapes)
        except Exception:
            pass

p = Presentation(PATH)
for si, slide in enumerate(p.slides):
    print(f"===== SLIDE {si} =====")
    for sh in iter_shapes(slide.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip():
            print("  TXT:", repr(sh.text_frame.text))
        if sh.has_table:
            for row in sh.table.rows:
                print("  TBL:", [c.text for c in row.cells])
